import streamlit as st
import sympy as sp
import numpy as np
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import schemdraw
import schemdraw.elements as elm
import io
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# 1. Konfigurasi Halaman & Tema Light Mode
st.set_page_config(layout="wide", page_title="Analisis Transien Rangkaian Listrik", page_icon="⚡")
import matplotlib
matplotlib.use('Agg')
schemdraw.use('matplotlib')

# Fungsi-fungsi SchemDraw (Menyimpan sebagai PNG untuk Support PPTX & Light Theme)
def draw_full_schematic():
    path = "schema_full.svg"
    with schemdraw.Drawing(show=False) as d:
        d.config(color='black', bgcolor='white', lw=2)
        V1 = d.add(elm.SourceSin().up().label('$v(t)$', loc='left'))
        R1 = d.add(elm.Resistor().right().label('$4\\,\\Omega$'))
        d.push()
        Rsh = d.add(elm.Resistor().down().label('$8\\,\\Omega$', loc='bottom'))
        L1 = d.add(elm.Inductor().left().label('$2\\text{ H}$', loc='bottom'))
        
        d.pop()
        R2 = d.add(elm.Resistor().right().label('$8\\,\\Omega$'))
        Line2 = d.add(elm.Line().down())
        L2 = d.add(elm.Inductor().left().label('$4\\text{ H}$', loc='bottom'))
        
        c1x = (V1.start[0] + Rsh.start[0]) / 2
        c1y = (V1.start[1] + R1.end[1]) / 2
        d += elm.LoopArrow(width=1.2, height=1.2, direction='cw').at((c1x, c1y)).label('$i_1$', color='#00a2ed').color('#00a2ed')
        
        c2x = (Rsh.start[0] + Line2.end[0]) / 2
        c2y = (R2.start[1] + L2.start[1]) / 2
        d += elm.LoopArrow(width=1.2, height=1.2, direction='cw').at((c2x, c2y)).label('$i_2$', color='#00a2ed').color('#00a2ed')
        try:
            d.save(path)
        except Exception:
            pass
    return path

def draw_loop1_only():
    path = "schema_loop1.svg"
    with schemdraw.Drawing(show=False) as d:
        d.config(color='black', bgcolor='white', lw=2)
        V1 = d.add(elm.SourceSin().up().label('$v(t)$', loc='left'))
        R1 = d.add(elm.Resistor().right().label('$4\\,\\Omega$'))
        d.push()
        Rsh = d.add(elm.Resistor().down().label('$8\\,\\Omega$', loc='bottom'))
        L1 = d.add(elm.Inductor().left().label('$2\\text{ H}$', loc='bottom'))
        
        d.pop()
        d.add(elm.Line().right().length(1))
        d.add(elm.Dot(open=True))
        
        c1x = (V1.start[0] + Rsh.start[0]) / 2
        c1y = (V1.start[1] + R1.end[1]) / 2
        d += elm.LoopArrow(width=1.2, height=1.2, direction='cw').at((c1x, c1y)).label('$i_1$', color='#00a2ed').color('#00a2ed')
        try:
            d.save(path)
        except Exception:
            pass
    return path

def draw_loop2_only():
    path = "schema_loop2.svg"
    with schemdraw.Drawing(show=False) as d:
        d.config(color='black', bgcolor='white', lw=2)
        d.add(elm.Dot(open=True))
        d.add(elm.Line().right().length(1))
        d.push()
        Rsh = d.add(elm.Resistor().down().label('$8\\,\\Omega$', loc='bottom'))
        d.add(elm.Line().left().tox(d.elements[0].start))
        
        d.pop()
        R2 = d.add(elm.Resistor().right().label('$8\\,\\Omega$'))
        Line2 = d.add(elm.Line().down())
        L2 = d.add(elm.Inductor().left().label('$4\\text{ H}$', loc='bottom'))
        
        c2x = (Rsh.start[0] + Line2.end[0]) / 2
        c2y = (R2.start[1] + L2.start[1]) / 2
        d += elm.LoopArrow(width=1.2, height=1.2, direction='cw').at((c2x, c2y)).label('$i_2$', color='#00a2ed').color('#00a2ed')
        try:
            d.save(path)
        except Exception:
            pass
    return path

def create_presentation(fig_arus):
    prs = Presentation()
    
    # Slide 1: Judul
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Analisis Rangkaian Transien Dua Loop"
    slide.placeholders[1].text = "Parameter:\nv(t) = 390cos(t)\ni(0) = 0, i'(0) = 0\nLoop 1: L1=2H, R1=4\u03A9, Rsh=8\u03A9\nLoop 2: L2=4H, R2=8\u03A9"
    
    # Slide 2: Skematik
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Skematik Rangkaian (Loop 1 & Loop 2)"
    try:
        if os.path.exists("schema_loop1.png"):
            slide.shapes.add_picture("schema_loop1.png", Inches(0.5), Inches(2.5), width=Inches(4))
        if os.path.exists("schema_loop2.png"):
            slide.shapes.add_picture("schema_loop2.png", Inches(5), Inches(2.5), width=Inches(4))
    except Exception as e:
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
        txBox.text_frame.text = f"Gambar skematik gagal dimuat: {e}"
        
    # Slide 3: Kesimpulan Soal 1
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Penyelesaian Soal 1 (Sumber Kontinu)"
    tf = slide.placeholders[1].text_frame
    tf.text = "Fungsi Arus dalam Domain Waktu:"
    p = tf.add_paragraph()
    p.text = "i1(t) = -26 \u00B7 exp(-2t) - 16 \u00B7 exp(-8t) + 42 \u00B7 cos(t) + 15 \u00B7 sin(t)"
    p = tf.add_paragraph()
    p.text = "i2(t) = -26 \u00B7 exp(-2t) + 8 \u00B7 exp(-8t) + 18 \u00B7 cos(t) + 12 \u00B7 sin(t)"
    
    # Slide 4: Kesimpulan Soal 2
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Penyelesaian Soal 2 (Sumber Terputus di t=\u03C0)"
    tf = slide.placeholders[1].text_frame
    tf.text = "Dengan Time-Shifting Laplace Transform, forced response menghilang di t \u2265 \u03C0:"
    p = tf.add_paragraph()
    p.text = "i1_putus(t) = i1(t) + i1(t-\u03C0)u(t-\u03C0)"
    p = tf.add_paragraph()
    p.text = "i2_putus(t) = i2(t) + i2(t-\u03C0)u(t-\u03C0)"
    
    # Slide 5: Grafik Respons
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Dinamika Arus (Soal 1 & Soal 2)"
    try:
        img_bytes = fig_arus.to_image(format="png", width=900, height=450)
        img_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(img_stream, Inches(0.5), Inches(2), width=Inches(9))
    except Exception as e:
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
        txBox.text_frame.text = f"Gagal mengekspor grafik Plotly.\nPastikan library 'kaleido' terinstall: pip install -U kaleido\nError detail: {e}"
        
    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

# --- UI Setup ---
st.title("⚡ Laporan Teknis Digital: Analisis Rangkaian Listrik Transien")
st.markdown("---")

def render_latex(formula):
    st.latex(r"\displaystyle " + formula)


# --- 1. Eksposisi Masalah ---
st.header("1. Eksposisi Masalah")
col_exp1, col_exp2 = st.columns([1.5, 1])
with col_exp1:
    st.markdown("""
    Kita akan menganalisis sebuah rangkaian dua loop dengan parameter-parameter berikut:
    - **Tegangan Sumber**: $v(t) = 390 \\cos(t)$ V
    - **Kondisi Awal**: $i(0) = 0$, $i'(0) = 0$
    - **Loop 1**: $R_1 = 4\\,\\Omega$, $L_1 = 2\\text{ H}$, dan resistor bersama $R_{shared} = 8\\,\\Omega$
    - **Loop 2**: $R_2 = 8\\,\\Omega$, $L_2 = 4\\text{ H}$, dan resistor bersama $R_{shared} = 8\\,\\Omega$

    Analisis ini dilakukan menggunakan **Transformasi Laplace** untuk mengekstraksi respons arus secara menyeluruh, mencakup *natural response* maupun *forced response*.
    """)
with col_exp2:
    try:
        path_full = draw_full_schematic()
        with open(path_full, "r") as f:
            st.markdown(f'<div style="text-align:center; padding:10px; background:white; border-radius:10px;">{f.read()}</div>', unsafe_allow_html=True)
    except:
        st.warning("Skematik utama gagal dimuat.")
st.markdown("---")

# --- 2. Penurunan Matematis (Deep Dive) ---
st.header("2. Penurunan Matematis Langkah-Demi-Langkah")

# Setup Symbolic Variables SymPy
t, s = sp.symbols('t s')
I1, I2 = sp.symbols('I_1 I_2')

st.markdown("Pendekatan pedagogis kita memecah rangkaian menjadi dua loop independen untuk mempermudah analisis KVL awal:")

col_L1, col_L2 = st.columns(2)
with col_L1:
    st.markdown("**Analisis Loop 1 (Kiri)**")
    try:
        path_l1 = draw_loop1_only()
        with open(path_l1, "r") as f:
            st.markdown(f'<div style="text-align:center; padding:10px; background:white; border-radius:10px;">{f.read()}</div>', unsafe_allow_html=True)
    except:
        st.warning("Gagal memuat skema Loop 1.")
    render_latex(r"L_1 \frac{di_1}{dt} + R_1 i_1 + R_{sh}(i_1 - i_2) = v(t)")
    render_latex(r"2 \frac{di_1}{dt} + 12i_1 - 8i_2 = v(t)")

with col_L2:
    st.markdown("**Analisis Loop 2 (Kanan)**")
    try:
        path_l2 = draw_loop2_only()
        with open(path_l2, "r") as f:
            st.markdown(f'<div style="text-align:center; padding:10px; background:white; border-radius:10px;">{f.read()}</div>', unsafe_allow_html=True)
    except:
        st.warning("Gagal memuat skema Loop 2.")
    render_latex(r"L_2 \frac{di_2}{dt} + R_2 i_2 + R_{sh}(i_2 - i_1) = 0")
    render_latex(r"4 \frac{di_2}{dt} + 16i_2 - 8i_1 = 0")

st.markdown("---")
# ==========================================
# SOAL 1
# ==========================================
st.subheader("Penyelesaian Soal 1: Sumber Tegangan Kontinu")
st.markdown(r"Pada soal pertama, asumsi tegangan eksitasi tidak pernah terputus, yaitu $v(t) = 390\cos(t)$. Transformasi Laplace untuk sistem ini menghasilkan:")

v_t1 = 390 * sp.cos(t)
V_s1 = sp.laplace_transform(v_t1, t, s, noconds=True).simplify()

render_latex(r"\mathcal{L}\{v(t)\} = " + sp.latex(V_s1))
render_latex(r"\begin{bmatrix} 2s + 12 & -8 \\ -8 & 4s + 16 \end{bmatrix} \begin{bmatrix} I_1(s) \\ I_2(s) \end{bmatrix} = \begin{bmatrix} " + sp.latex(V_s1) + r" \\ 0 \end{bmatrix}")

mat_A = sp.Matrix([[2*s + 12, -8], [-8, 4*s + 16]])
mat_B = sp.Matrix([[V_s1], [0]])
sol = mat_A.LUsolve(mat_B)

I1_s = sp.factor(sol[0])
I2_s = sp.factor(sol[1])

st.markdown("Isolasi variabel arus pada domain-s (Laplace):")
render_latex(r"I_1(s) = " + sp.latex(I1_s))
render_latex(r"I_2(s) = " + sp.latex(I2_s))

st.markdown("### Dekomposisi Pecahan Parsial untuk $I_2(s)$:")
render_latex(r"I_2(s) = \frac{390s}{(s+2)(s+8)(s^2+1)} = \frac{A}{s+2} + \frac{B}{s+8} + \frac{Cs+D}{s^2+1}")
st.markdown("Mencari A dan B (Metode Residu):")
render_latex(r"A = \left. \frac{390s}{(s+8)(s^2+1)} \right|_{s=-2} = \frac{-780}{(6)(5)} = -26")
render_latex(r"B = \left. \frac{390s}{(s+2)(s^2+1)} \right|_{s=-8} = \frac{-3120}{(-6)(65)} = 8")
st.markdown("Mencari C dan D (Substitusi Nilai):")
st.markdown("Substitusi $s=0$:")
render_latex(r"0 = \frac{-26}{2} + \frac{8}{8} + D \implies D = 12")
st.markdown("Substitusi $s=1$:")
render_latex(r"\frac{390(1)}{(3)(9)(2)} = \frac{-26}{3} + \frac{8}{9} + \frac{C+12}{2} \implies \frac{65}{9} = \frac{C+12}{2} - \frac{70}{9} \implies C = 18")
render_latex(r"I_2(s) = \frac{-26}{s+2} + \frac{8}{s+8} + \frac{18s+12}{s^2+1}")
render_latex(r"i_{2,Soal1}(t) = -26e^{-2t} + 8e^{-8t} + 18\cos(t) + 12\sin(t)")

st.markdown("### Dekomposisi Pecahan Parsial untuk $I_1(s)$:")
render_latex(r"I_1(s) = \frac{195s(s+4)}{(s+2)(s+8)(s^2+1)} = \frac{A}{s+2} + \frac{B}{s+8} + \frac{Cs+D}{s^2+1}")
st.markdown("Mencari A dan B (Metode Residu):")
render_latex(r"A = \left. \frac{195s(s+4)}{(s+8)(s^2+1)} \right|_{s=-2} = \frac{-780}{30} = -26")
render_latex(r"B = \left. \frac{195s(s+4)}{(s+2)(s^2+1)} \right|_{s=-8} = \frac{6240}{-390} = -16")
st.markdown("Mencari C dan D (Substitusi Nilai):")
st.markdown("Substitusi $s=0$:")
render_latex(r"0 = \frac{-26}{2} - \frac{16}{8} + D \implies D = 15")
st.markdown("Substitusi $s=1$:")
render_latex(r"\frac{975}{54} = \frac{-26}{3} - \frac{16}{9} + \frac{C+15}{2} \implies C = 42")
render_latex(r"I_1(s) = \frac{-26}{s+2} - \frac{16}{s+8} + \frac{42s+15}{s^2+1}")
render_latex(r"i_{1,Soal1}(t) = -26e^{-2t} - 16e^{-8t} + 42\cos(t) + 15\sin(t)")

# ==========================================
# SOAL 2
# ==========================================
st.markdown("---")
st.subheader(r"Penyelesaian Soal 2: Sumber Tegangan Terputus di $t=3\pi$")
st.markdown(r"Tegangan eksitasi untuk soal 2 dimodelkan dengan Heaviside Step Function $u(t)$:")
render_latex(r"v_2(t) = 390\cos(t)[u(t) - u(t-3\pi)]")
st.markdown(r"Menggunakan sifat trigonometri $\cos(t-3\pi) = -\cos(t)$:")
render_latex(r"v_2(t) = 390\cos(t)u(t) + 390\cos(t-3\pi)u(t-3\pi)")

st.markdown(r"### 1. Transformasi Laplace & Teorema Time-Shifting:")
render_latex(r"V_2(s) = \frac{390s}{s^2+1} + e^{-3\pi s}\frac{390s}{s^2+1} = \frac{390s}{s^2+1}(1+e^{-3\pi s})")
st.markdown(r"Respons arus s-domain secara otomatis memuat pergeseran $(1+e^{-3\pi s})$ terhadap Soal 1:")
render_latex(r"I_{1,Soal2}(s) = I_1(s) + I_1(s)e^{-3\pi s}")
render_latex(r"I_{2,Soal2}(s) = I_2(s) + I_2(s)e^{-3\pi s}")

st.markdown(r"### 2. Invers Laplace ke Ranah Waktu:")
render_latex(r"i_{1,Soal2}(t) = i_{1,Soal1}(t)u(t) + i_{1,Soal1}(t-3\pi)u(t-3\pi)")
render_latex(r"i_{2,Soal2}(t) = i_{2,Soal1}(t)u(t) + i_{2,Soal1}(t-3\pi)u(t-3\pi)")

st.markdown(r"### 3. Evaluasi Dinamika Sistem untuk $t>3\pi$:")
st.markdown(r"Pada rentang $t>3\pi$, nilai step $u(t)=1$ dan $u(t-3\pi)=1$. Sifat pergeseran fasa: $\sin(t-3\pi)=-\sin(t)$ dan $\cos(t-3\pi)=-\cos(t)$.")
st.markdown(r"Penjabaran Arus $i_1(t)$ pada Soal 2 untuk $t>3\pi$:")
render_latex(r"i_{1,Soal2}(t) = [15\sin(t) + 42\cos(t) - 16e^{-8t} - 26e^{-2t}] + [-15\sin(t) - 42\cos(t) - 16e^{-8(t-3\pi)} - 26e^{-2(t-3\pi)}]")
st.markdown(r"Komponen steady-state saling membatalkan:")
render_latex(r"i_{1,Soal2}(t) = -16e^{-8t}(1+e^{24\pi}) - 26e^{-2t}(1+e^{6\pi}) \text{ Ampere}")

st.markdown(r"Penjabaran Arus $i_2(t)$ pada Soal 2 untuk $t>3\pi$:")
render_latex(r"i_{2,Soal2}(t) = [12\sin(t) + 18\cos(t) + 8e^{-8t} - 26e^{-2t}] + [-12\sin(t) - 18\cos(t) + 8e^{-8(t-3\pi)} - 26e^{-2(t-3\pi)}]")
st.markdown(r"Komponen steady-state saling membatalkan:")
render_latex(r"i_{2,Soal2}(t) = 8e^{-8t}(1+e^{24\pi}) - 26e^{-2t}(1+e^{6\pi}) \text{ Ampere}")

st.markdown("---")

# --- 3. Analisis Stabilitas: Pole-Zero Plot ---
st.header("3. Analisis Stabilitas: Pole-Zero Plot")
col_pz1, col_pz2 = st.columns([1.5, 1])

with col_pz1:
    fig_pz = go.Figure()
    # Poles: dari penyebut (s+2)(s+8)(s^2+1) -> -2, -8, +j, -j
    poles_real = [-2, -8, 0, 0]
    poles_imag = [0, 0, 1, -1]
    # Zeros for I1(s) -> s=0, s=-4
    zeros_real = [0, -4]
    zeros_imag = [0, 0]
    
    fig_pz.add_trace(go.Scatter(x=poles_real, y=poles_imag, mode='markers', 
                                marker=dict(symbol='x', color='red', size=12, line_width=2), 
                                name='Poles'))
    fig_pz.add_trace(go.Scatter(x=zeros_real, y=zeros_imag, mode='markers', 
                                marker=dict(symbol='circle-open', color='blue', size=10, line_width=2), 
                                name='Zeros (I1)'))
    
    fig_pz.add_vline(x=0, line_dash="solid", line_color="black", line_width=1.5)
    fig_pz.add_hline(y=0, line_dash="solid", line_color="black", line_width=1.5)
    
    fig_pz.update_layout(
        title="Pole-Zero Map (s-plane)",
        xaxis_title="Real (\u03C3)",
        yaxis_title="Imaginary (j\u03C9)",
        template="plotly_white",
        width=600, height=400,
        xaxis=dict(range=[-10, 2]),
        yaxis=dict(range=[-2, 2])
    )
    st.plotly_chart(fig_pz, width="stretch")

with col_pz2:
    st.markdown(r"""
    **Analisis Stabilitas Sistem:**
    - **Poles di Sumbu Imajiner ($\pm j$ atau $+j$ dan $-j$)**: Berasal dari fungsi eksitasi input tegangan. Poles ini merepresentasikan respons osilasi *steady-state* yang konstan.
    - **Poles di Sumbu Real Negatif (-2 dan -8)**: Menunjukkan karakteristik intrinsik rangkaian (R dan L). Letaknya berada di *Left-Half Plane* (LHP), yang mengartikan bahwa sistem bersifat **stabil secara asimtotik**. Tanpa paksaan dari luar, respons alamiahnya akan berupa peluruhan eksponensial ($e^{-2t}$ dan $e^{-8t}$) menuju angka nol murni.
    """)

st.markdown("---")

# --- 4. Analisis Respons Transien & Visualisasi ---
st.header("4. Analisis Respons Transien")

# Perpanjang waktu hingga 25 detik
t_vals = np.linspace(0, 25, 1000)

# Array Soal 1: Kontinu
i1_cont = -26*np.exp(-2*t_vals) - 16*np.exp(-8*t_vals) + 42*np.cos(t_vals) + 15*np.sin(t_vals)
i2_cont = -26*np.exp(-2*t_vals) + 8*np.exp(-8*t_vals) + 18*np.cos(t_vals) + 12*np.sin(t_vals)

# Array Soal 2: Putus pada t=3pi
i1_putus = np.where(t_vals < 3*np.pi, 
                    i1_cont, 
                    -26*np.exp(-2*t_vals)*(1 + np.exp(6*np.pi)) - 16*np.exp(-8*t_vals)*(1 + np.exp(24*np.pi)))
i2_putus = np.where(t_vals < 3*np.pi,
                    i2_cont,
                    -26*np.exp(-2*t_vals)*(1 + np.exp(6*np.pi)) + 8*np.exp(-8*t_vals)*(1 + np.exp(24*np.pi)))

fig_arus = make_subplots(rows=1, cols=2, subplot_titles=("Arus Loop 1 (i1)", "Arus Loop 2 (i2)"), shared_yaxes=True)

# Plot Loop 1
fig_arus.add_trace(go.Scatter(x=t_vals, y=i1_cont, mode='lines', name='Soal 1 (i1 Kontinu)', line=dict(color='#1f77b4', width=2)), row=1, col=1)
fig_arus.add_trace(go.Scatter(x=t_vals, y=i1_putus, mode='lines', name='Soal 2 (i1 Putus t=π)', line=dict(color='#ff7f0e', dash='dash', width=2.5)), row=1, col=1)

# Plot Loop 2
fig_arus.add_trace(go.Scatter(x=t_vals, y=i2_cont, mode='lines', name='Soal 1 (i2 Kontinu)', line=dict(color='#2ca02c', width=2)), row=1, col=2)
fig_arus.add_trace(go.Scatter(x=t_vals, y=i2_putus, mode='lines', name='Soal 2 (i2 Putus t=π)', line=dict(color='#d62728', dash='dash', width=2.5)), row=1, col=2)

# Garis putus pada t = pi
pi_val = 3 * np.pi
fig_arus.add_vline(x=pi_val, line_dash="dot", line_color="black", row=1, col=1)
fig_arus.add_vline(x=pi_val, line_dash="dot", line_color="black", row=1, col=2)

fig_arus.update_layout(
    template="plotly_white", 
    title_text="Dinamika Respons Arus Rangkaian (Soal 1 & Soal 2 Side-by-Side)",
    hovermode="x unified",
    legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1)
)
fig_arus.update_xaxes(title_text="Waktu (detik)", row=1, col=1)
fig_arus.update_xaxes(title_text="Waktu (detik)", row=1, col=2)
fig_arus.update_yaxes(title_text="Arus (Ampere)", row=1, col=1)

st.plotly_chart(fig_arus, width="stretch")

st.markdown(r"""
### Narasi Analisis Fisis
**Observasi Grafik:**
1. **Soal 1 (Garis Solid)**: Pada detik awal ($0 \leq t < 3$ detik), terjadi efek *transien* di mana amplitudo belum stabil dan gelombang terdistorsi oleh redaman eksponensial. Setelah melampaui rentang $t > 5$ detik, arus sepenuhnya masuk ke dalam wujud *steady-state* osilasi murni tanpa redaman akibat suplai paksa berkelanjutan dari tegangan sumber.
2. **Soal 2 (Garis Putus-Putus)**: Tepat pada $t = 3\pi \approx 9.42$ detik (ditandai garis vertikal), interupsi eksitasi terjadi. Secara visual, arus pada grafik tidak jatuh lurus menghantam 0 Ampere, membuktikan inersia akibat induktor $L$. Secara empiris, *forced response* lenyap, menyisakan *natural response* peluruhan eksponensial murni yang meluruh ke arah ekuilibrium. Saat waktu diekspansi hingga $25$ detik, energi magnetik di dalam kumparan dipastikan telah terkuras seutuhnya (menjadi panas oleh resistor $R$).
""")

# --- Export PPTX Section ---
st.markdown("---")
st.subheader("Ekspor ke Presentasi")
if HAS_PPTX:
    try:
        ppt_file = create_presentation(fig_arus)
        st.download_button(
            label="Download Laporan PPTX 📥",
            data=ppt_file,
            file_name="Laporan_Analisis_Transien.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    except Exception as e:
        st.error(f"Gagal memproses PPTX: {e}")
else:
    st.warning("Library `python-pptx` belum terpasang.")
    st.info("Untuk mengaktifkan fitur Export ke PPTX, jalankan perintah di terminal:\n\n`pip install python-pptx kaleido`")
